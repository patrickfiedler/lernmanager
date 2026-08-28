"""Deterministic artifact gate checks — no LLM, no external services.

Called on upload when a subtask has artifact_gate_json set.
Returns {passed, message, details, matches} — student sees both what's missing
and what already looks good.
"""

import io
import zipfile
import json
from difflib import SequenceMatcher


def check_gate(file_bytes: bytes, filename: str, gate_config: dict) -> dict:
    """Dispatch to format-specific check. Passes by default for unknown formats."""
    ext = ('.' + filename.rsplit('.', 1)[-1].lower()) if '.' in filename else ''
    if ext in ('.pptx', '.odp'):
        return _check_presentation(file_bytes, ext, gate_config)
    if ext in ('.docx', '.odt'):
        return _check_document(file_bytes, ext, gate_config)
    if ext == '.sb3':
        return _check_scratch(file_bytes, gate_config)
    return {'passed': True, 'message': '', 'details': [], 'matches': []}


def _fuzzy_match(a: str, b: str) -> float:
    return SequenceMatcher(None, a.lower().strip(), b.lower().strip()).ratio()


def check_filename(filename: str, expected_filename: str, student_vorname: str = '', student_name: str = '') -> dict:
    """Deterministic filename check — never sent to the LLM (see conventions.md).

    Compares the uploaded filename's stem against expected_filename with
    [Vorname]/[Name] placeholders substituted for the student's real name.
    Shaped like an LLM checklist item so it can be spliced into the same
    feedback list; the caller marks it non-LLM (e.g. source='deterministic').
    """
    stem = filename.rsplit('.', 1)[0] if '.' in filename else filename
    ext = filename[len(stem):]
    expected = expected_filename.replace('[Vorname]', student_vorname).replace('[Name]', student_name)
    passed = stem.strip().lower() == expected.strip().lower()
    # Show the name the student should end up with, extension included -- they
    # compare against a whole filename, not a stem, so spelling out "ohne
    # Dateiendung" only made the message harder to act on. expected_filename is
    # authored as a stem (conventions.md), but tolerate one that already has it.
    if ext and not expected.lower().endswith(ext.lower()):
        expected_display = expected + ext
    else:
        expected_display = expected
    # criterion carries the requirement, note says what the student's file
    # actually looks like -- same split as every LLM criterion. Repeating the
    # target name in both just said the same thing twice.
    note = (
        "Der Dateiname ist korrekt."
        if passed
        else f'Deine Datei heißt „{filename}".'
    )
    return {'criterion': f'Dateiname ist „{expected_display}"', 'passed': passed,
            'note': note, 'source': 'deterministic'}


def _result(issues: list, matches: list = None, warnings: list = None) -> dict:
    passed = not issues
    message = "Abgabe sieht vollständig aus ✓" if passed else "Abgabe noch nicht vollständig"
    return {'passed': passed, 'message': message, 'details': issues, 'matches': matches or [], 'warnings': warnings or []}


# --- text rules: required_text / forbidden_text / expect_content_in ---
#
# These read the blocks extract_*_blocks() returns, so they can ask *where* a
# string appears, not only whether it appears. Scoping is optional:
#
#   required_text: ["Alle fünf Fachraumregeln"]        # anywhere
#   required_text: [{text: "Fachraumregeln", kind: heading}]
#   forbidden_text: [{text: "[Dein Name]", in: header}]

# Region names as an author would write them, mapped to block regions.
_RULE_REGIONS = {
    'body': 'body', 'slide': 'slide', 'slides': 'slide', 'notes': 'notes',
    'header': 'header', 'footer': 'footer',
}

_REGION_LABELS = {
    'body': 'im Dokument', 'slide': 'auf den Folien', 'notes': 'in den Notizen',
    'header': 'in der Kopfzeile', 'footer': 'in der Fußzeile',
}


def _text_rule(entry):
    """Normalize one rule. Accepts a bare string or {text, in, kind}."""
    if isinstance(entry, str) and entry.strip():
        return {'text': entry}
    if isinstance(entry, dict) and isinstance(entry.get('text'), str) and entry['text'].strip():
        return entry
    return None


def _scoped_blocks(blocks: list, rule: dict) -> list:
    region = _RULE_REGIONS.get(str(rule.get('in', '')).strip().lower())
    kind = rule.get('kind')
    if region:
        blocks = [b for b in blocks if b['region'] == region]
    if kind:
        blocks = [b for b in blocks if b['kind'] == kind]
    return blocks


def _text_present(blocks: list, rule: dict, threshold: float) -> bool:
    """True when the rule's text is in scope somewhere.

    Substring, case-insensitive. That is what forbidden_text lives on -- a
    template placeholder like "______" sits inside a longer line ("Name: ______
    Klasse: ____"), so a whole-line comparison would never see it.

    Fuzzy whole-block matching applies to kind: heading only, where it preserves
    what required_headings did. It is deliberately not offered for body text:
    at the default threshold of 0.6, "Kein Essen am PC." scores 0.74 against
    "Keine Getränke am PC.", so a document holding one Fachraumregel would pass
    a check for a different one -- exactly the check MBI needs this field for.
    Headings are few and distinct; body lines are many and similar.
    """
    needle = ' '.join(rule['text'].split()).lower()
    fuzzy = rule.get('kind') == 'heading'
    for b in _scoped_blocks(blocks, rule):
        text = ' '.join(b['text'].split()).lower()
        if needle in text:
            return True
        if fuzzy and _fuzzy_match(needle, text) >= threshold:
            return True
    return False


def _rule_label(rule: dict) -> str:
    """What to call the thing in student-facing feedback."""
    return 'Abschnitt' if rule.get('kind') == 'heading' else 'Text'


def _check_text_rules(blocks: list, config: dict, issues: list, matches: list, warnings: list):
    """Apply required_text, forbidden_text and expect_content_in to blocks."""
    threshold = config.get('title_match_threshold', 0.6)

    for entry in config.get('required_text', []):
        rule = _text_rule(entry)
        if rule is None:
            continue
        label = _rule_label(rule)
        if _text_present(blocks, rule, threshold):
            matches.append(f'{label} gefunden: „{rule["text"]}" ✓')
        else:
            issues.append(f'{label} fehlt: „{rule["text"]}"')

    for entry in config.get('forbidden_text', []):
        rule = _text_rule(entry)
        if rule is None:
            continue
        if _text_present(blocks, rule, threshold):
            issues.append(f'Noch aus der Vorlage übrig: „{rule["text"]}"')
        else:
            matches.append(f'„{rule["text"]}" ist ersetzt ✓')

    # A warning, never a failure -- the "wrote it in the speaker notes" case.
    expect = _RULE_REGIONS.get(str(config.get('expect_content_in', '')).strip().lower())
    if expect:
        total = sum(len(b['text'].split()) for b in blocks)
        inside = sum(len(b['text'].split()) for b in blocks if b['region'] == expect)
        if total and inside / total < 0.5:
            warnings.append(f"Der meiste Text steht nicht {_REGION_LABELS[expect]}")


def _check_presentation(file_bytes: bytes, ext: str, config: dict) -> dict:
    """Check slide count, required titles (fuzzy), min chars per slide, and min images.

    config keys:
      format (list[str]) — accepted extensions, checked by caller
      min_slides (int)
      min_images (int)
      required_slide_titles (list[str])
      required_text / forbidden_text / expect_content_in — see _check_text_rules
      title_match_threshold (float) — fuzzy ratio, default 0.6
      min_chars_per_slide (int)
    """
    issues = []
    matches = []
    warnings = []
    threshold = config.get('title_match_threshold', 0.6)

    if ext == '.odp':
        import xml.etree.ElementTree as ET
        NS = {
            'draw': 'urn:oasis:names:tc:opendocument:xmlns:drawing:1.0',
            'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0',
            'presentation': 'urn:oasis:names:tc:opendocument:xmlns:presentation:1.0',
        }
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                root = ET.fromstring(z.read('content.xml'))
        except Exception:
            return {'passed': False, 'message': 'Datei konnte nicht gelesen werden', 'details': ['Ungültige .odp-Datei'], 'matches': []}

        pages = root.findall('.//draw:page', NS)
        slide_count = len(pages)
        if config.get('min_slides', 0):
            if slide_count < config['min_slides']:
                issues.append(f"Zu wenig Folien ({slide_count}, erwartet: {config['min_slides']})")
            else:
                matches.append(f"{slide_count} Folien ✓")

        titles = []
        for page in pages:
            title_text = ''
            for frame in page.findall('.//draw:frame', NS):
                cls = frame.get('{urn:oasis:names:tc:opendocument:xmlns:presentation:1.0}class')
                if cls == 'title':
                    title_text = ' '.join(frame.itertext()).strip()
                    break
            titles.append(title_text)

        for req in config.get('required_slide_titles', []):
            if max((_fuzzy_match(req, t) for t in titles), default=0) < threshold:
                issues.append(f'Folie fehlt: „{req}"')
            else:
                matches.append(f'Folie gefunden: „{req}" ✓')

        min_chars = config.get('min_chars_per_slide', 0)
        if min_chars:
            for i, page in enumerate(pages, 1):
                text = ' '.join(page.itertext()).strip()
                if len(text) < min_chars:
                    issues.append(f"Folie {i} hat zu wenig Text ({len(text)} Zeichen, erwartet: {min_chars})")

    elif ext == '.pptx':
        from pptx import Presentation
        try:
            prs = Presentation(io.BytesIO(file_bytes))
        except Exception:
            return {'passed': False, 'message': 'Datei konnte nicht gelesen werden', 'details': ['Ungültige .pptx-Datei'], 'matches': []}

        slides = prs.slides
        slide_count = len(slides)
        if config.get('min_slides', 0):
            if slide_count < config['min_slides']:
                issues.append(f"Zu wenig Folien ({slide_count}, erwartet: {config['min_slides']})")
            else:
                matches.append(f"{slide_count} Folien ✓")

        titles = [slide.shapes.title.text if slide.shapes.title else '' for slide in slides]
        for req in config.get('required_slide_titles', []):
            if max((_fuzzy_match(req, t) for t in titles), default=0) < threshold:
                issues.append(f'Folie fehlt: „{req}"')
            else:
                matches.append(f'Folie gefunden: „{req}" ✓')

        min_chars = config.get('min_chars_per_slide', 0)
        if min_chars:
            for i, slide in enumerate(slides, 1):
                text = ' '.join(s.text for s in slide.shapes if hasattr(s, 'text')).strip()
                if len(text) < min_chars:
                    issues.append(f"Folie {i} hat zu wenig Text ({len(text)} Zeichen, erwartet: {min_chars})")

    min_images = config.get('min_images', 0)
    if min_images:
        image_exts = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.wmf', '.emf', '.svg'}
        prefix = 'Pictures/' if ext == '.odp' else 'ppt/media/'
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                image_count = sum(
                    1 for name in z.namelist()
                    if name.startswith(prefix) and ('.' + name.rsplit('.', 1)[-1].lower()) in image_exts
                )
        except Exception:
            image_count = 0
        if image_count < min_images:
            issues.append(f"Zu wenig Bilder ({image_count}, erwartet: {min_images})")
        else:
            matches.append(f"{image_count} Bild{'er' if image_count != 1 else ''} ✓")

    # Text rules read blocks, which carry the slide number and tell slide text
    # apart from speaker notes -- neither branch above can see that.
    import artifact_processor
    try:
        blocks = (artifact_processor.extract_pptx_blocks(file_bytes) if ext == '.pptx'
                  else artifact_processor.extract_odp_blocks(file_bytes))
    except Exception:
        blocks = []
    _check_text_rules(blocks, config, issues, matches, warnings)

    return _result(issues, matches, warnings)


def _check_document(file_bytes: bytes, ext: str, config: dict) -> dict:
    """Check required headings (fuzzy), minimum word count, and minimum image count.

    config keys:
      format (list[str]) — accepted extensions, checked by caller
      min_words (int)
      min_words_required (bool) — promote min_words from warning to failure
      min_images (int)
      required_text / forbidden_text / expect_content_in — see _check_text_rules
      title_match_threshold (float) — fuzzy ratio, default 0.6
    """
    import artifact_processor
    try:
        blocks = (artifact_processor.extract_docx_blocks(file_bytes) if ext == '.docx'
                  else artifact_processor.extract_odt_blocks(file_bytes))
    except Exception:
        return {'passed': False, 'message': 'Datei konnte nicht gelesen werden', 'details': ['Ungültige Datei']}

    issues = []
    matches = []
    warnings = []
    min_words = config.get('min_words', 0)
    if min_words:
        # Counted on the rendered string, not on the block texts: the '#'
        # heading markers count as words there, and every min_words value MBI
        # authored so far was tuned against that. Fixing the wart belongs with
        # min_added_words, not here.
        word_count = len(artifact_processor.render_text(blocks).split())
        if word_count < min_words:
            # min_words is a warning by default -- a short but finished artifact
            # should not be blocked. min_words_required is for the from-scratch
            # ones, where too little text really does mean unfinished.
            if config.get('min_words_required'):
                issues.append(f"Zu wenig Text ({word_count} Wörter, erwartet: {min_words})")
            else:
                warnings.append("wenig Text vorhanden")
        else:
            matches.append("Wortanzahl erreicht")

    min_images = config.get('min_images', 0)
    if min_images:
        image_exts = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.wmf', '.emf', '.svg'}
        prefix = 'word/media/' if ext == '.docx' else 'Pictures/'
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                image_count = sum(
                    1 for name in z.namelist()
                    if name.startswith(prefix) and ('.' + name.rsplit('.', 1)[-1].lower()) in image_exts
                )
        except Exception:
            image_count = 0
        if image_count < min_images:
            issues.append(f"Zu wenig Bilder ({image_count}, erwartet: {min_images})")
        else:
            matches.append(f"{image_count} Bild{'er' if image_count != 1 else ''} ✓")

    # required_headings used to live here, parsing '#' prefixes back out of the
    # flat string. It is now required_text with kind: heading -- one field for
    # one operation, and it reaches body text too.
    _check_text_rules(blocks, config, issues, matches, warnings)

    return _result(issues, matches, warnings)


def _check_scratch(file_bytes: bytes, config: dict) -> dict:
    """Check sprite count and script count in a .sb3 project."""
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            project = json.loads(z.read('project.json'))
    except Exception:
        return {'passed': False, 'message': 'Datei konnte nicht gelesen werden', 'details': ['Ungültige .sb3-Datei']}

    issues = []
    matches = []
    sprites = [t for t in project.get('targets', []) if not t.get('isStage')]

    if config.get('min_sprites', 0):
        if len(sprites) < config['min_sprites']:
            issues.append(f"Zu wenig Figuren ({len(sprites)}, erwartet: {config['min_sprites']})")
        else:
            matches.append(f"{len(sprites)} Figuren ✓")

    script_count = sum(
        sum(1 for b in s.get('blocks', {}).values() if isinstance(b, dict) and b.get('topLevel'))
        for s in sprites
    )
    if config.get('min_scripts', 0):
        if script_count < config['min_scripts']:
            issues.append(f"Zu wenig Skripte ({script_count}, erwartet: {config['min_scripts']})")
        else:
            matches.append(f"{script_count} Skripte ✓")

    return _result(issues, matches)
