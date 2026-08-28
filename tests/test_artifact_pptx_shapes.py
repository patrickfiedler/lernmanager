"""A .pptx table and a grouped shape are slide content too.

has_text_frame is False for a table (python-pptx models it as a GraphicFrame)
and a flat loop over slide.shapes does not descend into a group, so both were
dropped. .odp's blind `.//text:p` sweep caught them, which is how the two
presentation formats disagreed on the same deck.
"""
import io

import pytest
from pptx import Presentation
from pptx.util import Inches

from artifact_checker import _check_presentation
from artifact_processor import extract_pptx_blocks


def deck_with(table=False, group=False) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(4), Inches(1))
    box.text_frame.text = "Auf der Folie"
    if table:
        t = slide.shapes.add_table(1, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
        t.cell(0, 0).text = "Keine Getränke am PC"
        t.cell(0, 1).text = "Flüssigkeit zerstört die Tastatur"
    if group:
        g = slide.shapes.add_group_shape()
        inner = g.shapes.add_textbox(Inches(1), Inches(4), Inches(3), Inches(1))
        inner.text_frame.text = "In der Gruppe"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def texts(blocks):
    return [b['text'] for b in blocks]


def test_table_cells_are_extracted_and_tagged_as_cells():
    blocks = extract_pptx_blocks(deck_with(table=True))

    assert "Keine Getränke am PC" in texts(blocks)
    assert [b['kind'] for b in blocks if b['text'] == "Keine Getränke am PC"] == ['table-cell']


def test_text_inside_a_group_is_extracted():
    blocks = extract_pptx_blocks(deck_with(group=True))

    assert "In der Gruppe" in texts(blocks)


def test_a_required_text_rule_can_find_a_table_cell():
    res = _check_presentation(deck_with(table=True), '.pptx',
                              {'required_text': [{'text': "Keine Getränke am PC",
                                                  'kind': 'table-cell'}]})

    assert res['passed'] is True


def test_table_text_counts_toward_min_chars_per_slide():
    without = _check_presentation(deck_with(), '.pptx', {'min_chars_per_slide': 40})
    with_table = _check_presentation(deck_with(table=True), '.pptx', {'min_chars_per_slide': 40})

    assert without['passed'] is False
    assert with_table['passed'] is True


def test_everything_still_carries_the_right_slide_number():
    blocks = extract_pptx_blocks(deck_with(table=True, group=True))

    assert {b['index'] for b in blocks} == {1}
