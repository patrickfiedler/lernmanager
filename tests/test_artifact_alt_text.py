"""Alt text is extracted as kind: 'alt-text'.

Accessibility is a teaching point in MBI's units and alt text is ideal
deterministic-check material -- a fixed field with a fixed expectation. It was
invisible in every format: Word stores it on <wp:docPr descr>, PowerPoint on
<p:cNvPr descr>, ODF as <svg:desc> inside the frame.

Like images and headers, alt text does not join the rendered string: it was
never part of what the LLM graded.
"""
import base64
import io
import zipfile

from artifact_checker import _check_document, _check_presentation
from artifact_processor import (
    extract_docx_blocks,
    extract_odp_blocks,
    extract_pptx_blocks,
    render_text,
)

W = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
WP = 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'
A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_ODF = ('xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
        'xmlns:draw="urn:oasis:names:tc:opendocument:xmlns:drawing:1.0" '
        'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" '
        'xmlns:svg="urn:oasis:names:tc:opendocument:xmlns:svg-compatible:1.0" '
        'xmlns:xlink="http://www.w3.org/1999/xlink"')

ALT = "Foto einer Tastatur"

# 1x1 transparent PNG.
PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg==")


def alt_texts(blocks):
    return [b['text'] for b in blocks if b['kind'] == 'alt-text']


def docx_with_alt() -> bytes:
    body = (f'<w:p><w:r><w:t>Fließtext</w:t></w:r></w:p>'
            f'<w:p><w:r><w:drawing xmlns:wp="{WP}">'
            f'<wp:docPr id="1" name="Bild 1" descr="{ALT}"/>'
            f'<a:blip xmlns:a="{A}" xmlns:r="http://schemas.openxmlformats.org/'
            f'officeDocument/2006/relationships" r:embed="rId4"/>'
            f'</w:drawing></w:r></w:p>')
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w') as z:
        z.writestr('word/document.xml',
                   f'<w:document xmlns:w="{W}"><w:body>{body}</w:body></w:document>')
    return buf.getvalue()


def odp_with_alt() -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w') as z:
        z.writestr('content.xml',
                   f'<office:document-content {_ODF}><office:body><office:presentation>'
                   '<draw:page>'
                   '<draw:frame presentation:class="title" '
                   'xmlns:presentation="urn:oasis:names:tc:opendocument:xmlns:presentation:1.0">'
                   '<text:p>Meine Geräte</text:p></draw:frame>'
                   f'<draw:frame><svg:desc>{ALT}</svg:desc>'
                   '<draw:image xlink:href="Pictures/foto.jpg"/></draw:frame>'
                   '</draw:page></office:presentation></office:body></office:document-content>')
    return buf.getvalue()


def pptx_with_alt() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(io.BytesIO(PNG), Inches(1), Inches(1), Inches(1), Inches(1))
    pic._element._nvXxPr.cNvPr.set('descr', ALT)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_docx_alt_text_is_extracted():
    assert alt_texts(extract_docx_blocks(docx_with_alt())) == [ALT]


def test_odp_alt_text_is_extracted_with_its_slide_number():
    blocks = [b for b in extract_odp_blocks(odp_with_alt()) if b['kind'] == 'alt-text']

    assert [(b['text'], b['index']) for b in blocks] == [(ALT, 1)]


def test_pptx_alt_text_is_extracted():
    assert ALT in alt_texts(extract_pptx_blocks(pptx_with_alt()))


def test_alt_text_stays_out_of_the_string_the_llm_grades():
    assert render_text(extract_docx_blocks(docx_with_alt())) == "Fließtext"
    assert render_text(extract_odp_blocks(odp_with_alt())) == "[Folie 1]\nMeine Geräte"


def test_a_gate_can_require_alt_text_on_a_picture():
    res = _check_document(docx_with_alt(), '.docx',
                          {'required_text': [{'text': "Tastatur", 'kind': 'alt-text'}]})

    assert res['passed'] is True


def test_a_missing_alt_text_is_reported():
    res = _check_presentation(odp_with_alt(), '.odp',
                              {'required_text': [{'text': "Bildschirm", 'kind': 'alt-text'}]})

    assert res['passed'] is False
    assert 'Text fehlt: „Bildschirm"' in res['details']
