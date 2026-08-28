"""Artifact text extraction and pseudonymization for LLM feedback.

All functions operate on bytes in memory — no temp files, no disk storage.
This satisfies the DSGVO requirement that original files never persist server-side.

Supported formats: .pptx, .odp, .docx, .odt, .sb3
"""

import io
import re
import zipfile
import xml.etree.ElementTree as ET



# --- Block contract ---
#
# Extractors report *blocks* -- one line of text plus where it came from -- so a
# check can ask "does this string appear in a heading" instead of only "does it
# appear somewhere". Everything the old flat-string extraction knew and threw
# away (list item, speaker note, slide number) survives here.
#
#   region: body | slide | notes | header | footer | alt-text | comment
#   kind:   heading | paragraph | list-item | table-cell | title | caption
#           | alt-text | image
#   index:  slide/page number, where the format has one
#   level:  1-6, on kind='heading' only
#
# Deliberately not modelled: nesting depth, style names, geometry, formatting.
# Nothing needs them, and each would be a surface to maintain.

REGIONS = ('body', 'slide', 'notes', 'header', 'footer', 'alt-text', 'comment')
KINDS = ('heading', 'paragraph', 'list-item', 'table-cell', 'title', 'caption',
         'alt-text', 'image')


def block(text: str, region: str = 'body', kind: str = 'paragraph',
          index: int = None, level: int = None) -> dict:
    """Build one block. Optional keys stay absent rather than None."""
    b = {'text': text, 'region': region, 'kind': kind}
    if index is not None:
        b['index'] = index
    if level is not None:
        b['level'] = level
    return b


# Reported for checks, never part of the flat string the LLM reads.
_UNRENDERED_KINDS = ('image', 'alt-text')


def _render_line(b: dict) -> str:
    """One block as one line of the flat string: headings keep their # markers."""
    if b['kind'] == 'heading':
        return f"{'#' * min(b.get('level', 1), 6)} {b['text']}"
    return b['text']


def render_text(blocks: list) -> str:
    """Render blocks back into the flat string the LLM prompt and the student's
    transparency view have always seen.

    This is the migration seam: every artifact grade so far was calibrated
    against this exact string, so it must not drift.

    Documents are one line per block. Presentations open each slide with a
    [Folie N] marker and are separated by a blank line; a slide that produced no
    blocks gets no marker, because the old extractor skipped it too.

    Blocks that are not lines of text -- a placed image, an alt-text caption --
    are reported by the extractors for the gate to count, but they were never
    part of this string and must not join it now.
    """
    sections = []
    lines = []
    current = None  # slide number of the section being built, None for body text
    for b in blocks:
        if b['kind'] in _UNRENDERED_KINDS or not b['text']:
            continue
        index = b.get('index')
        if index != current:
            if lines:
                sections.append('\n'.join(lines))
            lines = [f"[Folie {index}]"] if index is not None else []
            current = index
        lines.append(_render_line(b))
    if lines:
        sections.append('\n'.join(lines))
    return '\n\n'.join(sections)


# --- .pptx extraction ---

def _pptx_is_picture(shape) -> bool:
    """A shape that carries actual image bytes -- Picture, or a placeholder
    filled with one. Asking for .image is the reliable test; shape_type does not
    distinguish a filled picture placeholder from an empty one."""
    try:
        shape.image
        return True
    except (AttributeError, ValueError, KeyError):
        return False


def extract_pptx_blocks(file_bytes: bytes) -> list:
    """Blocks for a .pptx file, one per non-empty paragraph, in slide order."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    blocks = []
    for i, slide in enumerate(prs.slides, start=1):
        # python-pptx hands out a fresh proxy object on every .title access, so
        # compare the underlying XML element, not the wrapper.
        title = slide.shapes.title
        title_elem = title._element if title is not None else None
        for shape in slide.shapes:
            if _pptx_is_picture(shape):
                blocks.append(block('', 'slide', 'image', index=i))
                continue
            if not shape.has_text_frame:
                continue
            # kind says where the text came from, not what it means: every
            # paragraph of the title placeholder is a title block, even when the
            # author typed a whole slide's worth of lines into it.
            kind = 'title' if shape._element is title_elem else 'paragraph'
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    blocks.append(block(line, 'slide', kind, index=i))
    return blocks


def extract_pptx(file_bytes: bytes) -> str:
    """Extract slide text from a .pptx file. Returns one section per slide."""
    return render_text(extract_pptx_blocks(file_bytes))


def strip_pptx_metadata(file_bytes: bytes) -> bytes:
    """Remove author/creator/lastModifiedBy from .pptx core properties.

    .pptx files are ZIP archives. The core properties live in
    docProps/core.xml. We blank the relevant fields before any storage.
    """
    _CLEAR_TAGS = {
        '{http://purl.org/dc/elements/1.1/}creator',
        '{http://schemas.openxmlformats.org/package/2006/metadata/core-properties}lastModifiedBy',
    }
    buf = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(file_bytes), 'r') as zin, \
         zipfile.ZipFile(buf, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == 'docProps/core.xml':
                root = ET.fromstring(data)
                for elem in root.iter():
                    if elem.tag in _CLEAR_TAGS:
                        elem.text = ''
                data = ET.tostring(root, encoding='unicode').encode('utf-8')
            zout.writestr(item, data)
    return buf.getvalue()


# --- .odp extraction ---

_ODP_NS = {
    'text':         'urn:oasis:names:tc:opendocument:xmlns:text:1.0',
    'draw':         'urn:oasis:names:tc:opendocument:xmlns:drawing:1.0',
    'presentation': 'urn:oasis:names:tc:opendocument:xmlns:presentation:1.0',
    'office':       'urn:oasis:names:tc:opendocument:xmlns:office:1.0',
}


def _odp_collect_blocks(elem, blocks, index, region, kind):
    """Walk one <draw:page>, carrying region/kind context down the tree.

    A <presentation:notes> child switches the region -- everything inside it is
    speaker notes, not slide text. A <draw:frame presentation:class="title">
    switches the kind. The nearest enclosing container wins.
    """
    for child in elem:
        tag = child.tag.split('}')[-1]
        if tag == 'notes':
            _odp_collect_blocks(child, blocks, index, 'notes', 'paragraph')
        elif tag == 'frame':
            cls = child.get('{%(presentation)s}class' % _ODP_NS)
            _odp_collect_blocks(child, blocks, index, region,
                                'title' if cls == 'title' else kind)
        elif tag == 'list-item':
            _odp_collect_blocks(child, blocks, index, region, 'list-item')
        elif tag == 'image':
            # A <draw:image> is an image actually placed on the slide. Counting
            # ZIP entries under Pictures/ instead counts orphans: this very
            # template ships a 128 KB JPEG referenced only from manifest.xml.
            blocks.append(block('', region, 'image', index=index))
        elif tag == 'p':
            text = _odf_line_text(child).strip()
            if text:
                blocks.append(block(text, region, kind, index=index))
        else:
            _odp_collect_blocks(child, blocks, index, region, kind)


def extract_odp_blocks(file_bytes: bytes) -> list:
    """Blocks for an .odp file (ODF Presentation).

    .odp is a ZIP archive. Slides are <draw:page> elements in content.xml;
    text sits under <draw:frame> → <draw:text-box> → <text:p>.
    """
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        root = ET.fromstring(z.read('content.xml'))
    blocks = []
    for i, page in enumerate(root.findall('.//{%(draw)s}page' % _ODP_NS), start=1):
        _odp_collect_blocks(page, blocks, i, 'slide', 'paragraph')
    return blocks


def extract_odp(file_bytes: bytes) -> str:
    """Extract slide text from an .odp file (ODF Presentation)."""
    return render_text(extract_odp_blocks(file_bytes))


# --- Pseudonymization ---

def anonymize(text: str, student_name: str, class_name: str = '') -> str:
    """Replace student name and class name occurrences with placeholders.

    student_name should be "Vorname Nachname" or just one name.
    Handles first name, last name, and combined forms separately.
    class_name (e.g. "5a") is replaced with [Klasse].

    If student_name is empty, name replacement is skipped.
    """
    parts = student_name.split()
    if parts:
        # Build patterns: each individual name part + both combined orders
        candidates = list(parts)
        if len(parts) >= 2:
            candidates.append(r'\s+'.join(re.escape(p) for p in parts))           # First Last
            candidates.append(r'\s+'.join(re.escape(p) for p in reversed(parts))) # Last First

        # Apply longest patterns first to avoid partial replacements
        for pattern in sorted(candidates, key=len, reverse=True):
            escaped = pattern if r'\s+' in pattern else re.escape(pattern)
            text = re.sub(r'\b' + escaped + r'\b', '[Schüler/in]', text, flags=re.IGNORECASE)

    if class_name.strip():
        for pattern in _class_name_patterns(class_name.strip()):
            text = re.sub(pattern, '[Klasse]', text, flags=re.IGNORECASE)

    return text


def _class_name_patterns(class_name: str) -> list[str]:
    """Build regex patterns for a class name, including typographic variants.

    Handles:
    - Exact match:                  "Ginkgo-Haie-Urvögel 5"
    - Hyphens/spaces interchangeable: "Ginkgo Haie Urvögel 5"
    - Without trailing grade number:  "Ginkgo-Haie-Urvögel" / "Ginkgo Haie Urvögel"
    """
    # Split on hyphens and whitespace to get tokens
    tokens = re.split(r'[\s\-]+', class_name)
    # Trailing token is the grade number if it's purely numeric
    has_grade = tokens and tokens[-1].isdigit()
    name_tokens = tokens[:-1] if has_grade else tokens
    grade_token = tokens[-1] if has_grade else None

    sep = r'[\s\-]+'  # flexible separator

    def joined(parts, include_grade=True):
        t = list(parts) + ([grade_token] if include_grade and grade_token else [])
        return r'\b' + sep.join(re.escape(p) for p in t) + r'\b'

    patterns = []
    if name_tokens:
        patterns.append(joined(name_tokens, include_grade=True))   # full name + grade
        if has_grade:
            patterns.append(joined(name_tokens, include_grade=False))  # name without grade
        # Individual compound parts (skip grade number — too common)
        for token in name_tokens:
            patterns.append(r'\b' + re.escape(token) + r'\b')
    else:
        patterns.append(r'\b' + re.escape(class_name) + r'\b')

    return patterns


# --- .docx extraction ---

_W = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'


def extract_docx_blocks(file_bytes: bytes) -> list:
    """Blocks for a .docx document.

    .docx is a ZIP archive. All body text lives in word/document.xml as <w:p>
    paragraphs -- including the ones inside table cells, which is why a flat
    iter() over <w:p> already reaches everything. Headings are identified by
    <w:pStyle w:val="Heading1"> (English Word) or "berschrift1" (German Word,
    the Ü is stripped in the XML); 'Title'/'Titel' is what add_heading(level=0)
    sets and carries no digit, so it falls back to level 1.
    """
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        root = ET.fromstring(z.read('word/document.xml'))
    # A paragraph does not know its ancestors, so build the reverse map once --
    # needed to tell a table cell's paragraph from a body paragraph.
    parents = {child: parent for parent in root.iter() for child in parent}
    blocks = []
    for para in root.iter(f'{{{_W}}}p'):
        pstyle = para.find(f'{{{_W}}}pPr/{{{_W}}}pStyle')
        style_val = pstyle.get(f'{{{_W}}}val', '') if pstyle is not None else ''
        text = ''.join(t.text or '' for t in para.iter(f'{{{_W}}}t')).strip()
        if not text:
            continue
        sl = style_val.lower()
        if sl.startswith('heading') or sl.startswith('berschrift') or sl in ('title', 'titel'):
            level = min(int(''.join(c for c in style_val if c.isdigit()) or '1'), 6)
            blocks.append(block(text, 'body', 'heading', level=level))
        else:
            blocks.append(block(text, 'body', _docx_kind(para, parents)))
    blocks.extend(_docx_image_blocks(root))
    return blocks


# DrawingML puts the picture reference in <a:blip>, the older VML in
# <v:imagedata>. Both mean an image actually placed in the document, unlike a
# stray file under word/media/.
_A_BLIP = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
_V_IMAGEDATA = '{urn:schemas-microsoft-com:vml}imagedata'


def _docx_image_blocks(root) -> list:
    return [block('', 'body', 'image')
            for e in root.iter() if e.tag in (_A_BLIP, _V_IMAGEDATA)]


def _docx_kind(para, parents) -> str:
    """Nearest enclosing container decides: numbering makes a list item, and a
    <w:tc> ancestor makes a table cell."""
    if para.find(f'{{{_W}}}pPr/{{{_W}}}numPr') is not None:
        return 'list-item'
    node = parents.get(para)
    while node is not None:
        if node.tag == f'{{{_W}}}tc':
            return 'table-cell'
        node = parents.get(node)
    return 'paragraph'


def extract_docx(file_bytes: bytes) -> str:
    """Extract text from a .docx document, preserving heading structure."""
    return render_text(extract_docx_blocks(file_bytes))


# --- .odt extraction ---

_ODF_TEXT_NS = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_ODF_OFFICE_NS = 'urn:oasis:names:tc:opendocument:xmlns:office:1.0'

# ODF wraps text in containers instead of keeping one flat paragraph stream the
# way .docx does. Only the leaf text:h / text:p elements carry text, so a walk
# that stops at the direct children of <office:text> loses every list and table.
# The value is the kind a paragraph inside that container gets.
_ODF_CONTAINERS = {
    'list': None, 'list-header': None, 'section': None,
    'table': None, 'table-row': None, 'table-header-rows': None,
    'list-item': 'list-item',
    'table-cell': 'table-cell',
}

# LibreOffice's "Titel" is a styled paragraph, not a text:h element -- the same
# trap as Word's Title style handled in extract_docx_blocks() above.
_ODF_TITLE_STYLES = {'title', 'titel'}


def _odf_line_text(elem) -> str:
    """Text of one ODF paragraph, expanding the elements ODF uses for whitespace.

    <text:s text:c="3"/> is a run of three spaces and <text:tab/> a tab.
    itertext() yields nothing for either, which collapsed the fill-in lines
    ("Name: ____   Klasse: ____") and made them differ from their .docx twin.
    """
    parts = [elem.text or '']
    for child in elem:
        tag = child.tag.split('}')[-1]
        if tag == 's':
            parts.append(' ' * int(child.get('{%s}c' % _ODF_TEXT_NS, '1') or 1))
        elif tag == 'tab':
            parts.append('\t')
        else:
            parts.append(_odf_line_text(child))
        parts.append(child.tail or '')
    return ''.join(parts)


def _odt_collect_blocks(elem, blocks, kind):
    """Walk ODF body content, descending into the containers that hold paragraphs."""
    for child in elem:
        tag = child.tag.split('}')[-1]
        if tag in _ODF_CONTAINERS:
            _odt_collect_blocks(child, blocks, _ODF_CONTAINERS[tag] or kind)
            continue
        if tag not in ('h', 'p'):
            continue
        text = _odf_line_text(child).strip()
        if not text:
            continue
        if tag == 'h':
            level = min(int(child.get('{%s}outline-level' % _ODF_TEXT_NS, '1') or 1), 6)
            blocks.append(block(text, 'body', 'heading', level=level))
        elif (child.get('{%s}style-name' % _ODF_TEXT_NS) or '').lower() in _ODF_TITLE_STYLES:
            blocks.append(block(text, 'body', 'heading', level=1))
        else:
            blocks.append(block(text, 'body', kind))


def extract_odt_blocks(file_bytes: bytes) -> list:
    """Blocks for an .odt document.

    .odt is a ZIP archive. Text lives in content.xml under <office:text>.
    Headings use <text:h text:outline-level="N">, paragraphs use <text:p>,
    and both can sit inside lists or table cells (see _ODF_CONTAINERS).
    """
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        root = ET.fromstring(z.read('content.xml'))
    body = root.find('.//{%s}text' % _ODF_OFFICE_NS)
    blocks = []
    if body is not None:
        _odt_collect_blocks(body, blocks, 'paragraph')
        # Images hang off <draw:frame> anchored inside a paragraph, so they sit
        # below the leaves the walk above stops at. Position carries no meaning
        # for a document, only the count does.
        blocks.extend(block('', 'body', 'image')
                      for e in body.iter() if e.tag.split('}')[-1] == 'image')
    return blocks


def extract_odt(file_bytes: bytes) -> str:
    """Extract text from an .odt document, preserving heading structure."""
    return render_text(extract_odt_blocks(file_bytes))


# --- .sb3 (Scratch) extraction ---

# Maps Scratch opcode prefixes to readable German category names
_SB3_CATEGORIES = {
    'motion':     'Bewegung',
    'looks':      'Aussehen',
    'sound':      'Klang',
    'event':      'Ereignisse',
    'control':    'Steuerung',
    'sensing':    'Fühlen',
    'operators':  'Operatoren',
    'data':       'Variablen',
    'procedures': 'Eigene Blöcke',
    'pen':        'Malstift',
}


def _sb3_collect_opcodes(blocks: dict) -> list[str]:
    """Return a list of all unique opcodes used in a target's blocks dict."""
    seen = set()
    result = []
    for block in blocks.values():
        if not isinstance(block, dict):
            continue
        opcode = block.get('opcode', '')
        if opcode and opcode not in seen:
            seen.add(opcode)
            result.append(opcode)
    return result



def _sb3_count_scripts(blocks: dict) -> int:
    """Count top-level blocks (= scripts) in a target's blocks dict."""
    return sum(1 for b in blocks.values() if isinstance(b, dict) and b.get('topLevel', False))


def _sb3_variable_names(target: dict) -> list[str]:
    """Return variable names defined on a target (stage = global, sprite = local)."""
    return [v[0] for v in target.get('variables', {}).values() if isinstance(v, list)]


def extract_sb3(file_bytes: bytes, filename: str = '') -> str:
    """Extract a readable project summary from a Scratch .sb3 file.

    .sb3 is a ZIP archive containing project.json. We collect per-target info
    (sprite name, opcode categories, costume/sound counts, variables, script count)
    and format it as structured text for LLM checklist evaluation.

    Level 1 criteria (block presence/counts) are fully supported.
    Level 2 criteria (script ordering, pseudo-code) require a future upgrade.
    """
    import json
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        project = json.loads(z.read('project.json'))

    targets = project.get('targets', [])
    stage = next((t for t in targets if t.get('isStage', False)), None)
    sprites = [t for t in targets if not t.get('isStage', False)]
    extensions = project.get('extensions', [])

    # Header counts
    backgrounds = len(stage.get('costumes', [])) if stage else 0
    global_vars = _sb3_variable_names(stage) if stage else []
    total_scripts = sum(_sb3_count_scripts(t.get('blocks', {})) for t in targets)
    total_blocks = sum(len([b for b in t.get('blocks', {}).values() if isinstance(b, dict)])
                       for t in targets)

    header_lines = ["[Scratch-Projekt]"]
    if filename:
        header_lines.append(f"Dateiname: {filename}")
    header_lines.append(
        f"Figuren: {len(sprites)} | Skripte: {total_scripts} | "
        f"Blöcke: {total_blocks} | Hintergründe: {backgrounds}"
    )
    if global_vars:
        header_lines.append(f"Globale Variablen: {', '.join(global_vars)}")
    header_lines.append(
        f"Erweiterungen: {', '.join(extensions) if extensions else '(keine)'}"
    )
    sections = ["\n".join(header_lines)]

    for target in targets:
        name = target.get('name', 'Unbekannt')
        is_stage = target.get('isStage', False)
        costumes = len(target.get('costumes', []))
        sounds = len(target.get('sounds', []))
        opcodes = _sb3_collect_opcodes(target.get('blocks', {}))
        local_vars = _sb3_variable_names(target) if not is_stage else []
        scripts = _sb3_count_scripts(target.get('blocks', {}))

        label = "Bühne" if is_stage else f"Figur: {name}"
        used = [_SB3_CATEGORIES[p] for p in _SB3_CATEGORIES
                if any(op.startswith(p + '_') for op in opcodes)]

        lines = [f"[{label}]"]
        if is_stage:
            lines.append(f"Hintergründe: {costumes}, Töne: {sounds}")
        else:
            lines.append(f"Skripte: {scripts}, Kostüme: {costumes}, Töne: {sounds}")
        lines.append(f"Kategorien: {', '.join(used) if used else '(keine)'}")
        if local_vars:
            lines.append(f"Lokale Variablen: {', '.join(local_vars)}")

        sections.append("\n".join(lines))

    return '\n\n'.join(sections)


# --- Format dispatch ---

ACCEPTED_FORMATS = {
    '.pptx': extract_pptx,
    '.odp':  extract_odp,
    '.docx': extract_docx,
    '.odt':  extract_odt,
    '.sb3':  extract_sb3,
}

_BLOCK_EXTRACTORS = {
    '.pptx': extract_pptx_blocks,
    '.odp':  extract_odp_blocks,
    '.docx': extract_docx_blocks,
    '.odt':  extract_odt_blocks,
}


def _artifact_ext(filename: str) -> str:
    return '.' + filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''


def extract_artifact_blocks(file_bytes: bytes, filename: str) -> list:
    """Blocks for a supported artifact file. Raises ValueError for unknown formats.

    .sb3 is the exception: a Scratch project has no text to read, so
    extract_sb3() writes a prose summary of the project instead. It comes back
    as a single block so callers never have to special-case the dispatch.
    """
    ext = _artifact_ext(filename)
    if ext not in ACCEPTED_FORMATS:
        raise ValueError(f"Unsupported format: {ext!r}. Accepted: {', '.join(ACCEPTED_FORMATS)}")
    if ext == '.sb3':
        return [block(extract_sb3(file_bytes, filename), 'body', 'paragraph')]
    return _BLOCK_EXTRACTORS[ext](file_bytes)


def extract_artifact(file_bytes: bytes, filename: str) -> str:
    """Extract text from a supported artifact file. Raises ValueError for unknown formats."""
    ext = _artifact_ext(filename)
    if ext == '.sb3':
        return extract_sb3(file_bytes, filename)
    return render_text(extract_artifact_blocks(file_bytes, filename))
